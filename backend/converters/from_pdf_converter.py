"""PDF转其他格式的转换器"""
import os
from typing import Optional
from .base_converter import BaseConverter


class PDFToWordConverter(BaseConverter):
    """PDF转Word转换器"""
    
    def convert(self, input_file: str, output_file: Optional[str] = None) -> str:
        """将PDF转换为Word"""
        try:
            from pdf2docx import Converter
        except ImportError:
            raise ImportError("需要安装pdf2docx库: pip install pdf2docx")
        
        self.validate_file(input_file, ['.pdf'])
        output_file = self.get_output_path(input_file, '.docx', output_file)
        
        try:
            cv = Converter(input_file)
            cv.convert(output_file)
            cv.close()
            return output_file
        except Exception as e:
            raise RuntimeError(f"PDF转Word失败: {str(e)}")


class PDFToPPTConverter(BaseConverter):
    """PDF转PowerPoint转换器"""
    
    def convert(self, input_file: str, output_file: Optional[str] = None) -> str:
        """将PDF转换为PowerPoint"""
        try:
            from pdf2image import convert_from_path
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            raise ImportError("需要安装pdf2image和python-pptx库: pip install pdf2image python-pptx")
        
        self.validate_file(input_file, ['.pdf'])
        output_file = self.get_output_path(input_file, '.pptx', output_file)
        
        try:
            # 将PDF转换为图片
            images = convert_from_path(input_file, dpi=200)
            
            # 创建PPT
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            for i, image in enumerate(images):
                # 保存临时图片
                temp_image = f"temp_page_{i}.png"
                image.save(temp_image, 'PNG')
                
                # 添加幻灯片
                blank_slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 添加图片到幻灯片
                left = top = Inches(0)
                slide.shapes.add_picture(temp_image, left, top, 
                                        width=prs.slide_width, 
                                        height=prs.slide_height)
                
                # 删除临时文件
                os.remove(temp_image)
            
            prs.save(output_file)
            return output_file
        except Exception as e:
            raise RuntimeError(f"PDF转PPT失败: {str(e)}\n提示: 需要安装poppler (Windows: https://github.com/oschwartz10612/poppler-windows/releases/)")


class PDFToImageConverter(BaseConverter):
    """PDF转图片转换器"""
    
    def convert(self, input_file: str, output_file: Optional[str] = None) -> str:
        """将PDF转换为图片（每页一张）"""
        try:
            from pdf2image import convert_from_path
        except ImportError:
            raise ImportError("需要安装pdf2image库: pip install pdf2image")
        
        self.validate_file(input_file, ['.pdf'])
        
        try:
            # 转换PDF为图片
            images = convert_from_path(input_file, dpi=300)
            
            # 生成输出文件名
            base_name = os.path.splitext(input_file)[0]
            output_files = []
            
            for i, image in enumerate(images, start=1):
                if output_file:
                    # 如果指定了输出文件，为多页添加序号
                    name_without_ext = os.path.splitext(output_file)[0]
                    ext = os.path.splitext(output_file)[1] or '.jpg'
                    page_output = f"{name_without_ext}_page_{i}{ext}"
                else:
                    page_output = f"{base_name}_page_{i}.jpg"
                
                image.save(page_output, 'JPEG')
                output_files.append(page_output)
            
            # 返回第一个文件路径或所有文件列表
            return output_files[0] if len(output_files) == 1 else ", ".join(output_files)
        except Exception as e:
            raise RuntimeError(f"PDF转图片失败: {str(e)}\n提示: 需要安装poppler")


class PDFToExcelConverter(BaseConverter):
    """PDF转Excel转换器"""
    
    def convert(self, input_file: str, output_file: Optional[str] = None) -> str:
        """将PDF转换为Excel"""
        try:
            import tabula
            import pandas as pd
        except ImportError:
            raise ImportError("需要安装tabula-py和pandas库: pip install tabula-py pandas openpyxl")
        
        self.validate_file(input_file, ['.pdf'])
        output_file = self.get_output_path(input_file, '.xlsx', output_file)
        
        try:
            # 读取PDF中的所有表格
            dfs = tabula.read_pdf(input_file, pages='all', multiple_tables=True)
            
            if not dfs:
                raise ValueError("PDF中未找到表格")
            
            # 写入Excel
            with pd.ExcelWriter(output_file, engine='openpyxl') as writer:
                for i, df in enumerate(dfs, start=1):
                    sheet_name = f'Sheet{i}'
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
            
            return output_file
        except Exception as e:
            raise RuntimeError(f"PDF转Excel失败: {str(e)}\n提示: 需要安装Java运行环境")
